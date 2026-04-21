# -*- coding: utf-8 -*-
"""
超无穹AI平台培训课程 - PPT文档生成
基于python-pptx库
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

OUTPUT_DIR = r"D:\网站部署\超无穹项目\chaowuqiong-project\docs\培训课程\文档生成\ppt"

def set_text_style(text_frame, font_size=18, bold=False, color=None, align='left'):
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = getattr(PP_ALIGN, align.upper())
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = RGBColor(*color)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(68, 114, 196)
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(220, 220, 220)

def add_section_slide(prs, section_title, section_num):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(112, 173, 71)
    background.line.fill.background()

    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Part {section_num}"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

def add_content_slide(prs, title, content_list, title_color=(68, 114, 196)):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*title_color)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.space_before = Pt(12)
        p.space_after = Pt(6)

def add_two_column_slide(prs, title, left_content, right_content, title_color=(68, 114, 196)):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*title_color)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(4.5), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_before = Pt(8)

    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_before = Pt(8)

def create_intro_ppt():
    """创建课程介绍PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "超无穹AI平台", "零基础小白培训课程")
    add_title_slide(prs, "培养全栈开发工程师", "12周从入门到就业")

    add_content_slide(prs, "课程概述", [
        "课程名称：超无穹AI平台全栈开发工程师培训",
        "目标学员：零基础小白，无需任何编程经验",
        "学习周期：12周（约3个月）",
        "总课时：72课时",
        "毕业水平：能够独立完成平台功能开发和部署",
    ])

    add_content_slide(prs, "学习路径", [
        "第一阶段（2周）：入门篇 - 计算机基础、HTML/CSS",
        "第二阶段（3周）：基础篇 - JavaScript、Node.js、数据库",
        "第三阶段（3周）：进阶篇 - React、平台功能开发",
        "第四阶段（2周）：高阶篇 - TypeScript、工程化、部署运维",
        "第五阶段（2周）：实战篇 - 综合项目、答辩、就业",
    ])

    add_content_slide(prs, "课程特色", [
        "Agent Orchestrator 智能辅助 - 多Agent协同自动分析学习进度",
        "渐进式学习曲线 - 从简单到复杂，步步为营",
        "实战项目驱动 - 每阶段都有实践项目，边学边做",
        "全栈能力培养 - 前端+后端+运维，全面发展",
        "企业级开发规范 - Git、代码规范、文档习惯",
    ])

    add_content_slide(prs, "实践项目", [
        "P1: 个人简历网站（HTML/CSS）⭐",
        "P2: 待办事项应用（JavaScript）⭐⭐",
        "P3: 用户管理系统（Node.js/MySQL）⭐⭐",
        "P4: AI对话组件（React）⭐⭐⭐",
        "P5-P9: 会员中心、支付功能、TS重构、Docker部署、综合项目",
    ])

    add_content_slide(prs, "毕业要求", [
        "出勤率 ≥ 90%（权重10%）",
        "作业完成：全部提交 + 及格以上（权重20%）",
        "阶段测试：每阶段≥60分（权重20%）",
        "实践项目：完成全部9个实践项目（权重30%）",
        "结业答辩：评委会≥60分通过（权重20%）",
    ])

    prs.save(os.path.join(OUTPUT_DIR, '01_课程介绍.pptx'))
    print(f"✅ 已生成：01_课程介绍.pptx")

def create_phase1_ppt():
    """创建第一阶段PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "第一阶段", "入门篇")

    add_section_slide(prs, "计算机基础与开发环境", "01")

    add_content_slide(prs, "第1周：计算机基础", [
        "课时1：计算机基本操作 - 操作系统使用、文件管理",
        "课时2：认识互联网 - URL、浏览器工作原理、网络基础",
        "课时3：编程是什么 - 编程思维、入门概念、学习方法",
        "课时4：开发工具介绍 - VS Code安装配置、终端使用",
        "课时5：超无穹平台介绍 - 平台功能演示、架构概述",
        "课时6：本周总结与答疑 - 知识点回顾、阶段测试",
    ], (68, 114, 196))

    add_content_slide(prs, "第2周：Web基础入门", [
        "课时7：HTML初识 - 标签、元素、属性、页面结构",
        "课时8：CSS入门 - 选择器、盒模型、常用样式",
        "课时9：响应式设计基础 - 媒体查询、Flexbox、栅格系统",
        "课时10：浏览器开发者工具 - Elements、Console、网络调试",
        "课时11：实践项目 - 综合运用HTML/CSS制作简历页面",
        "课时12：本周总结与答疑 - CSS常见问题、布局技巧",
    ], (68, 114, 196))

    add_two_column_slide(prs, "学习目标",
        ["理解计算机基本操作", "掌握文件管理技能", "熟悉浏览器使用", "建立编程思维"],
        ["掌握开发工具配置", "理解互联网基础", "掌握HTML标签", "掌握CSS样式"],
        (68, 114, 196))

    add_content_slide(prs, "实践任务", [
        "安装开发所需软件（VS Code、Git、Chrome）",
        "配置自己的开发环境",
        "编写个人介绍页面",
        "美化个人介绍页面样式",
        "实现页面自适应布局",
        "完成简历页面项目（P1）",
    ], (68, 114, 196))

    prs.save(os.path.join(OUTPUT_DIR, '02_第一阶段-入门篇.pptx'))
    print(f"✅ 已生成：02_第一阶段-入门篇.pptx")

def create_phase2_ppt():
    """创建第二阶段PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "第二阶段", "基础篇")

    add_section_slide(prs, "JavaScript基础", "02")

    add_content_slide(prs, "第3周：JavaScript基础（上）", [
        "课时13：JavaScript概述 - 发展历史、应用场景",
        "课时14：变量与数据类型 - let、const、数据类型转换",
        "课时15：运算符与表达式 - 算术、比较、逻辑运算符",
        "课时16：条件语句 - if/else、switch、三元运算符",
        "课时17：循环结构 - for、while、do-while",
        "课时18：函数基础 - 函数声明、调用、参数、返回值",
    ], (112, 173, 71))

    add_content_slide(prs, "第4周：JavaScript基础（下）", [
        "课时19：数组基础 - 创建、访问、遍历、常用方法",
        "课时20：对象基础 - 对象创建、属性访问、this关键字",
        "课时21：字符串处理 - 常用方法、正则表达式入门",
        "课时22：DOM操作 - 元素选择、创建、修改、事件绑定",
        "课时23：事件处理 - 事件类型、事件对象、事件委托",
        "课时24：实践项目 - 实现一个待办事项应用（Todo App）",
    ], (112, 173, 71))

    add_section_slide(prs, "Node.js与后端基础", "02")

    add_content_slide(prs, "第5周：Node.js与后端基础", [
        "课时25：Node.js介绍 - 概念、安装、npm包管理",
        "课时26：Express框架 - 路由、中间件、静态文件",
        "课时27：数据库概述 - MySQL安装、数据类型",
        "课时28：SQL基础 - CRUD操作、WHERE、ORDER BY",
        "课时29：API设计入门 - RESTful规范、接口文档",
        "课时30：阶段实践 - 完成用户注册登录功能（P3）",
    ], (112, 173, 71))

    add_two_column_slide(prs, "JavaScript核心技能",
        ["变量与数据类型", "运算符与表达式", "条件与循环语句", "函数与作用域"],
        ["数组与对象操作", "DOM操作", "事件处理", "ES6+新特性"],
        (112, 173, 71))

    prs.save(os.path.join(OUTPUT_DIR, '03_第二阶段-基础篇.pptx'))
    print(f"✅ 已生成：03_第二阶段-基础篇.pptx")

def create_phase3_ppt():
    """创建第三阶段PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "第三阶段", "进阶篇")

    add_section_slide(prs, "React框架入门", "03")

    add_content_slide(prs, "第6周：React框架入门", [
        "课时31：React概述 - 组件化、虚拟DOM、JSX",
        "课时32：组件与Props - 函数组件、Props传递、类型检查",
        "课时33：State与生命周期 - useState、useEffect",
        "课时34：事件处理 - 事件绑定、表单处理",
        "课时35：条件渲染与列表 - 列表渲染、key属性",
        "课时36：Hooks深入 - 自定义Hooks、Hooks规则",
    ], (237, 125, 49))

    add_section_slide(prs, "超无穹平台开发", "03")

    add_content_slide(prs, "第7周：超无穹平台前端开发", [
        "课时37：项目结构分析 - 前端目录、模块划分",
        "课时38：AI对话功能 - 对话UI、消息处理、流式响应",
        "课时39：角色管理 - 角色创建、编辑、删除功能",
        "课时40：会员系统前端 - 套餐展示、购买流程",
        "课时41：移动端适配 - 响应式设计、触摸交互",
        "课时42：实践项目 - 完成AI对话功能开发（P4）",
    ], (237, 125, 49))

    add_content_slide(prs, "第8周：后端API开发", [
        "课时43：项目结构分析 - 后端目录、路由划分",
        "课时44：认证中间件 - JWT、Token验证、权限控制",
        "课时45：LLM配置模块 - API代理、多模型切换",
        "课时46：支付集成 - 支付宝/微信支付、回调处理",
        "课时47：Redis缓存 - Session管理、缓存策略",
        "课时48：实践项目 - 完成对话历史API开发",
    ], (237, 125, 49))

    prs.save(os.path.join(OUTPUT_DIR, '04_第三阶段-进阶篇.pptx'))
    print(f"✅ 已生成：04_第三阶段-进阶篇.pptx")

def create_phase4_ppt():
    """创建第四阶段PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "第四阶段", "高阶篇")

    add_section_slide(prs, "TypeScript与工程化", "04")

    add_content_slide(prs, "第9周：TypeScript与工程化", [
        "课时49：TypeScript入门 - 类型系统、接口、泛型",
        "课时50：状态管理 - Redux/Zustand、状态设计模式",
        "课时51：前端路由 - React Router、动态路由、权限路由",
        "课时52：构建工具 - Vite打包优化、代码分割",
        "课时53：Git工作流 - 分支管理、代码合并、冲突解决",
        "课时54：代码规范 - ESLint、Prettier、提交规范",
    ], (112, 48, 160))

    add_section_slide(prs, "部署与运维", "04")

    add_content_slide(prs, "第10周：部署与运维", [
        "课时55：服务器基础 - Linux命令、SSH连接",
        "课时56：Nginx配置 - 反向代理、HTTPS配置",
        "课时57：PM2进程管理 - 进程守护、日志管理",
        "课时58：数据库运维 - 备份恢复、性能监控",
        "课时59：Docker基础 - 容器概念、镜像构建",
        "课时60：监控与日志 - 日志收集、异常告警",
    ], (112, 48, 160))

    add_two_column_slide(prs, "工程化技能",
        ["TypeScript类型系统", "状态管理模式", "前端路由设计", "性能优化技巧"],
        ["构建工具配置", "Git团队协作", "代码规范执行", "自动化测试"],
        (112, 48, 160))

    prs.save(os.path.join(OUTPUT_DIR, '05_第四阶段-高阶篇.pptx'))
    print(f"✅ 已生成：05_第四阶段-高阶篇.pptx")

def create_phase5_ppt():
    """创建第五阶段PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "第五阶段", "实战篇")

    add_section_slide(prs, "综合项目实战", "05")

    add_content_slide(prs, "第11周：综合项目实战", [
        "课时61：项目选题 - 团队组建、需求分析、技术选型",
        "课时62：系统设计 - 架构设计、数据库设计、接口设计",
        "课时63：前端开发 - 组件开发、页面构建、API对接",
        "课时64：后端开发 - API开发、数据库实现、业务逻辑",
        "课时65：测试与修复 - 单元测试、集成测试、Bug修复",
        "课时66：部署上线 - 生产环境部署、域名配置、SSL证书",
    ], (192, 0, 0))

    add_content_slide(prs, "第12周：结业与就业指导", [
        "课时67：项目答辩 - 项目展示、问题解答、评审打分",
        "课时68：课程总结 - 知识梳理、能力评估、成长回顾",
        "课时69：简历制作 - 技术栈梳理、项目经验整理",
        "课时70：面试指导 - 面试技巧、常见问题、薪资谈判",
        "课时71：职业规划 - 技术路线、行业发展、长期规划",
        "课时72：结业典礼 - 颁发证书、优秀表彰",
    ], (192, 0, 0))

    add_content_slide(prs, "毕业标准", [
        "出勤率 ≥ 90%（权重10%）",
        "作业完成：全部提交 + 及格以上（权重20%）",
        "阶段测试：每阶段≥60分（权重20%）",
        "实践项目：完成全部9个实践项目（权重30%）",
        "结业答辩：评委会≥60分通过（权重20%）",
    ], (192, 0, 0))

    add_title_slide(prs, "恭喜毕业！", "开启全栈开发工程师之路")

    prs.save(os.path.join(OUTPUT_DIR, '06_第五阶段-实战篇.pptx'))
    print(f"✅ 已生成：06_第五阶段-实战篇.pptx")

def main():
    print("=" * 50)
    print("超无穹AI平台培训课程 - PPT文档生成")
    print("=" * 50)

    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)

    create_intro_ppt()
    create_phase1_ppt()
    create_phase2_ppt()
    create_phase3_ppt()
    create_phase4_ppt()
    create_phase5_ppt()

    print("=" * 50)
    print(f"所有PPT文档已生成到：{OUTPUT_DIR}")
    print("=" * 50)

if __name__ == '__main__':
    main()
